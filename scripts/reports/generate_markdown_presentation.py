""" Generate PPTX Presentations From Repository Markdown Slide Decks """

from __future__ import annotations

# Import Python Utilities
import argparse, re
from dataclasses import dataclass
from pathlib import Path

# Import PPTX Utilities
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_PATH = Path(__file__).resolve().parents[2]
DEFAULT_TEMPLATE_PPTX_PATH = PROJECT_PATH / "reference" / "templates" / "Template_XiLab_Research.pptx"
MARKDOWN_COMMENT_PATTERN = re.compile(r"<!--.*?-->", re.DOTALL)

# Presentation Constants
TITLE_SLIDE_LAYOUT_INDEX = 0
CONTENT_SLIDE_LAYOUT_INDEX = 1

# Color Constants
TITLE_COLOR = RGBColor(22, 25, 59)
ACCENT_COLOR = RGBColor(53, 71, 140)
TEXT_COLOR = RGBColor(36, 43, 69)
TABLE_HEADER_FILL_COLOR = RGBColor(53, 71, 140)
TABLE_ROW_FILL_COLOR = RGBColor(250, 252, 255)
WHITE_COLOR = RGBColor(255, 255, 255)

@dataclass
class BulletBlock:

    """ Bullet Block """

    bullet_text_list: list[str]

@dataclass
class ParagraphBlock:

    """ Paragraph Block """

    paragraph_text_list: list[str]

@dataclass
class TableBlock:

    """ Table Block """

    header_cell_list: list[str]
    row_cell_matrix: list[list[str]]

@dataclass
class SlideContent:

    """ Slide Content """

    slide_title: str
    content_block_list: list[BulletBlock | ParagraphBlock | TableBlock]

def build_argument_parser() -> argparse.ArgumentParser:

    """ Build Argument Parser """

    argument_parser = argparse.ArgumentParser(
        description="Generate a repository-owned PPTX presentation from a Markdown slide deck."
    )

    # Configure Input And Output Paths
    argument_parser.add_argument("--input-markdown-path", required=True, help="Path to the Markdown slide-deck source.")
    argument_parser.add_argument("--output-pptx-path", default="", help="Optional explicit PPTX output path.")
    argument_parser.add_argument("--template-pptx-path", default="", help="Optional explicit PPTX template path.")

    return argument_parser

def parse_command_line_arguments() -> argparse.Namespace:

    """ Parse Command-Line Arguments """

    argument_parser = build_argument_parser()
    parsed_arguments = argument_parser.parse_args()
    return parsed_arguments

def resolve_input_markdown_path(input_markdown_path: str) -> Path:

    """ Resolve Input Markdown Path """

    resolved_input_markdown_path = Path(input_markdown_path).expanduser().resolve()
    assert resolved_input_markdown_path.exists(), f"Markdown deck path does not exist | {resolved_input_markdown_path}"
    assert resolved_input_markdown_path.is_file(), f"Markdown deck path is not a file | {resolved_input_markdown_path}"
    assert resolved_input_markdown_path.suffix.lower() == ".md", f"Markdown deck path is not a Markdown file | {resolved_input_markdown_path}"
    return resolved_input_markdown_path

def resolve_output_pptx_path(input_markdown_path: Path, output_pptx_path: str) -> Path:

    """ Resolve Output PPTX Path """

    if output_pptx_path:
        resolved_output_pptx_path = Path(output_pptx_path).expanduser().resolve()
    else:
        resolved_output_pptx_path = input_markdown_path.with_suffix(".pptx")

    assert resolved_output_pptx_path.suffix.lower() == ".pptx", f"Output path must end with .pptx | {resolved_output_pptx_path}"
    resolved_output_pptx_path.parent.mkdir(parents=True, exist_ok=True)
    return resolved_output_pptx_path

def resolve_template_pptx_path(template_pptx_path: str) -> Path:

    """ Resolve Template PPTX Path """

    resolved_template_pptx_path = Path(template_pptx_path).expanduser().resolve() if template_pptx_path else DEFAULT_TEMPLATE_PPTX_PATH.resolve()
    assert resolved_template_pptx_path.exists(), f"Template PPTX path does not exist | {resolved_template_pptx_path}"
    assert resolved_template_pptx_path.is_file(), f"Template PPTX path is not a file | {resolved_template_pptx_path}"
    assert resolved_template_pptx_path.suffix.lower() == ".pptx", f"Template PPTX path is not a PPTX file | {resolved_template_pptx_path}"
    return resolved_template_pptx_path

def load_markdown_text(input_markdown_path: Path) -> str:

    """ Load Markdown Text """

    markdown_text = input_markdown_path.read_text(encoding="utf-8")
    cleaned_markdown_text = MARKDOWN_COMMENT_PATTERN.sub("", markdown_text)
    return cleaned_markdown_text.strip()

def split_slide_segment_list(markdown_text: str) -> list[str]:

    """ Split Slide Segment List """

    slide_segment_list = [segment.strip() for segment in re.split(r"^\s*---\s*$", markdown_text, flags=re.MULTILINE) if segment.strip()]
    assert slide_segment_list, "The Markdown slide deck does not contain any slide segments."
    return slide_segment_list

def normalize_line_list(slide_segment_text: str) -> list[str]:

    """ Normalize Line List """

    return [line.rstrip() for line in slide_segment_text.splitlines()]

def split_block_text_list(slide_body_line_list: list[str]) -> list[list[str]]:

    """ Split Block Text List """

    block_text_list: list[list[str]] = []
    current_block_line_list: list[str] = []

    # Group Non-Empty Lines
    for current_line in slide_body_line_list:
        if current_line.strip():
            current_block_line_list.append(current_line)
            continue

        if current_block_line_list:
            block_text_list.append(current_block_line_list)
            current_block_line_list = []

    if current_block_line_list:
        block_text_list.append(current_block_line_list)

    return block_text_list

def parse_table_block(block_line_list: list[str]) -> TableBlock:

    """ Parse Table Block """

    stripped_line_list = [line.strip() for line in block_line_list]
    header_cell_list = [cell.strip() for cell in stripped_line_list[0].strip("|").split("|")]
    row_cell_matrix: list[list[str]] = []

    # Parse Table Rows
    for row_line in stripped_line_list[2:]:
        row_cell_list = [cell.strip() for cell in row_line.strip("|").split("|")]
        assert len(row_cell_list) == len(header_cell_list), "Markdown table row does not match header width."
        row_cell_matrix.append(row_cell_list)

    return TableBlock(
        header_cell_list=header_cell_list,
        row_cell_matrix=row_cell_matrix,
    )

def parse_content_block(block_line_list: list[str]) -> BulletBlock | ParagraphBlock | TableBlock:

    """ Parse Content Block """

    stripped_line_list = [line.strip() for line in block_line_list if line.strip()]
    assert stripped_line_list, "Encountered an empty content block."

    # Parse Markdown Table Block
    if all(line.startswith("|") and line.endswith("|") for line in stripped_line_list):
        assert len(stripped_line_list) >= 2, "Markdown table block is incomplete."
        return parse_table_block(stripped_line_list)

    # Parse Bullet Block
    if all(line.startswith("- ") for line in stripped_line_list):
        return BulletBlock(
            bullet_text_list=[line[2:].strip() for line in stripped_line_list],
        )

    # Parse Paragraph Block
    return ParagraphBlock(
        paragraph_text_list=[" ".join(stripped_line_list)],
    )

def sanitize_inline_markdown_text(text: str) -> str:

    """ Sanitize Inline Markdown Text """

    sanitized_text = text.strip()
    sanitized_text = re.sub(r"^#+\s+", "", sanitized_text)
    sanitized_text = sanitized_text.replace("`", "")
    sanitized_text = sanitized_text.replace("**", "")
    sanitized_text = sanitized_text.replace("*", "")
    sanitized_text = re.sub(r"\s+", " ", sanitized_text).strip()
    return sanitized_text

def parse_slide_segment(slide_segment_text: str, slide_index: int) -> SlideContent:

    """ Parse Slide Segment """

    normalized_line_list = normalize_line_list(slide_segment_text)
    non_empty_line_list = [line for line in normalized_line_list if line.strip()]
    assert non_empty_line_list, f"Slide segment is empty | slide_index={slide_index}"

    slide_heading_line = non_empty_line_list[0].strip()
    assert slide_heading_line.startswith("#"), f"Slide must start with a Markdown heading | slide_index={slide_index}"

    # Resolve Slide Title
    slide_title = slide_heading_line.lstrip("#").strip()
    assert slide_title, f"Slide heading is missing its title text | slide_index={slide_index}"

    # Resolve Body Lines
    heading_index = normalized_line_list.index(next(line for line in normalized_line_list if line.strip()))
    slide_body_line_list = normalized_line_list[heading_index + 1:]
    content_block_list = [
        parse_content_block(block_line_list)
        for block_line_list in split_block_text_list(slide_body_line_list)
    ]

    return SlideContent(
        slide_title=slide_title,
        content_block_list=content_block_list,
    )

def parse_slide_content_list(markdown_text: str) -> list[SlideContent]:

    """ Parse Slide Content List """

    slide_segment_list = split_slide_segment_list(markdown_text)
    slide_content_list = [
        parse_slide_segment(slide_segment_text, slide_index)
        for slide_index, slide_segment_text in enumerate(slide_segment_list, start=1)
    ]
    return slide_content_list

def create_presentation_from_template(template_pptx_path: Path) -> Presentation:

    """ Create Presentation From Template """

    presentation = Presentation(str(template_pptx_path))
    return presentation

def remove_slide(presentation: Presentation, slide_index: int) -> None:

    """ Remove Slide """

    slide_id_list = presentation.slides._sldIdLst
    slide_id = slide_id_list[slide_index]
    relationship_id = slide_id.rId
    presentation.part.drop_rel(relationship_id)
    slide_id_list.remove(slide_id)

def clear_template_slides(presentation: Presentation) -> None:

    """ Clear Template Slides """

    while len(presentation.slides) > 0:
        remove_slide(presentation, 0)

def resolve_slide_layout(presentation: Presentation, layout_index: int, fallback_index: int = 0):

    """ Resolve Slide Layout """

    if layout_index < len(presentation.slide_layouts):
        return presentation.slide_layouts[layout_index]

    return presentation.slide_layouts[fallback_index]

def resolve_slide_placeholder(slide, placeholder_index: int):

    """ Resolve Slide Placeholder """

    if placeholder_index < len(slide.placeholders):
        return slide.placeholders[placeholder_index]

    return None

def clear_text_frame(text_frame) -> None:

    """ Clear Text Frame """

    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

def resolve_content_bounds(slide) -> tuple[int, int, int, int]:

    """ Resolve Content Bounds """

    content_placeholder = resolve_slide_placeholder(slide, 1)

    if content_placeholder is not None:
        return (
            content_placeholder.left,
            content_placeholder.top,
            content_placeholder.width,
            content_placeholder.height,
        )

    return (
        Inches(1.0),
        Inches(1.55),
        Inches(11.0),
        Inches(5.4),
    )

def render_title_slide(presentation: Presentation, slide_content: SlideContent) -> None:

    """ Render Title Slide """

    slide_layout = resolve_slide_layout(presentation, TITLE_SLIDE_LAYOUT_INDEX)
    slide = presentation.slides.add_slide(slide_layout)

    # Style Title
    title_shape = slide.shapes.title
    title_shape.text = sanitize_inline_markdown_text(slide_content.slide_title)
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = TITLE_COLOR
    title_paragraph.alignment = PP_ALIGN.LEFT

    # Style Subtitle
    subtitle_shape = resolve_slide_placeholder(slide, 1)
    subtitle_text = ""

    if slide_content.content_block_list:
        first_block = slide_content.content_block_list[0]

        if isinstance(first_block, BulletBlock):
            subtitle_text = "\n".join(sanitize_inline_markdown_text(text) for text in first_block.bullet_text_list)
        elif isinstance(first_block, ParagraphBlock):
            subtitle_text = "\n".join(sanitize_inline_markdown_text(text) for text in first_block.paragraph_text_list)

    if subtitle_shape is not None:
        subtitle_shape.text = subtitle_text
        subtitle_text_frame = subtitle_shape.text_frame
        subtitle_text_frame.word_wrap = True

        for paragraph in subtitle_text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = ACCENT_COLOR
            paragraph.alignment = PP_ALIGN.LEFT

def render_text_blocks_into_placeholder(slide, content_block_list: list[BulletBlock | ParagraphBlock]) -> None:

    """ Render Text Blocks Into Placeholder """

    content_placeholder = resolve_slide_placeholder(slide, 1)

    if content_placeholder is None:
        content_placeholder = slide.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.0), Inches(5.4))

    text_frame = content_placeholder.text_frame
    clear_text_frame(text_frame)

    paragraph_is_first = True

    # Render Content Blocks
    for current_block in content_block_list:
        if isinstance(current_block, ParagraphBlock):
            for paragraph_text in current_block.paragraph_text_list:
                paragraph = text_frame.paragraphs[0] if paragraph_is_first else text_frame.add_paragraph()
                paragraph.text = sanitize_inline_markdown_text(paragraph_text)
                paragraph.font.size = Pt(22)
                paragraph.font.color.rgb = TEXT_COLOR
                paragraph.space_after = Pt(10)
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph_is_first = False
            continue

        if isinstance(current_block, BulletBlock):
            for bullet_text in current_block.bullet_text_list:
                paragraph = text_frame.paragraphs[0] if paragraph_is_first else text_frame.add_paragraph()
                paragraph.text = sanitize_inline_markdown_text(bullet_text)
                paragraph.level = 0
                paragraph.bullet = True
                paragraph.font.size = Pt(21)
                paragraph.font.color.rgb = TEXT_COLOR
                paragraph.space_after = Pt(8)
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph_is_first = False

def estimate_table_block_height(row_count: int) -> int:

    """ Estimate Table Block Height """

    return Inches(max(1.25, 0.55 + row_count * 0.42))

def add_table_block(slide, table_block: TableBlock) -> None:

    """ Add Table Block """

    column_count = len(table_block.header_cell_list)
    row_count = len(table_block.row_cell_matrix) + 1
    content_left, content_top, content_width, _ = resolve_content_bounds(slide)

    table_shape = slide.shapes.add_table(
        row_count,
        column_count,
        content_left,
        content_top,
        content_width,
        estimate_table_block_height(row_count),
    )
    table = table_shape.table

    # Populate Header Row
    for column_index, header_text in enumerate(table_block.header_cell_list):
        header_cell = table.cell(0, column_index)
        header_cell.text = sanitize_inline_markdown_text(header_text)
        header_cell.fill.solid()
        header_cell.fill.fore_color.rgb = TABLE_HEADER_FILL_COLOR
        header_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        header_paragraph = header_cell.text_frame.paragraphs[0]
        header_paragraph.font.bold = True
        header_paragraph.font.size = Pt(13)
        header_paragraph.font.color.rgb = WHITE_COLOR
        header_paragraph.alignment = PP_ALIGN.CENTER

    # Populate Data Rows
    for row_index, row_cell_list in enumerate(table_block.row_cell_matrix, start=1):
        for column_index, cell_text in enumerate(row_cell_list):
            table_cell = table.cell(row_index, column_index)
            table_cell.text = sanitize_inline_markdown_text(cell_text)
            table_cell.fill.solid()
            table_cell.fill.fore_color.rgb = TABLE_ROW_FILL_COLOR
            table_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell_paragraph = table_cell.text_frame.paragraphs[0]
            cell_paragraph.font.size = Pt(13)
            cell_paragraph.font.color.rgb = TEXT_COLOR
            cell_paragraph.alignment = PP_ALIGN.CENTER

def render_content_slide(presentation: Presentation, slide_content: SlideContent) -> None:

    """ Render Content Slide """

    slide_layout = resolve_slide_layout(presentation, CONTENT_SLIDE_LAYOUT_INDEX, fallback_index=TITLE_SLIDE_LAYOUT_INDEX)
    slide = presentation.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    if title_shape is not None:
        title_shape.text = sanitize_inline_markdown_text(slide_content.slide_title)
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.color.rgb = TITLE_COLOR
        title_paragraph.font.bold = True

    table_block_list = [current_block for current_block in slide_content.content_block_list if isinstance(current_block, TableBlock)]
    text_block_list = [current_block for current_block in slide_content.content_block_list if not isinstance(current_block, TableBlock)]

    if text_block_list:
        render_text_blocks_into_placeholder(slide, text_block_list)

    if table_block_list:
        add_table_block(slide, table_block_list[0])

def render_slide_content_list(presentation: Presentation, slide_content_list: list[SlideContent]) -> None:

    """ Render Slide Content List """

    assert slide_content_list, "At least one slide is required to render the presentation."

    # Render First Slide As Title Slide
    render_title_slide(presentation, slide_content_list[0])

    # Render Remaining Content Slides
    for slide_content in slide_content_list[1:]:
        render_content_slide(presentation, slide_content)

def print_generation_summary(input_markdown_path: Path, output_pptx_path: Path, template_pptx_path: Path, slide_count: int) -> None:

    """ Print Generation Summary """

    print("")
    print("================================================================================================")
    print("Presentation Generation Summary")
    print("================================================================================================")
    print(f"Input Markdown Path                {input_markdown_path}")
    print(f"Template PPTX Path                 {template_pptx_path}")
    print(f"Output PPTX Path                   {output_pptx_path}")
    print(f"Rendered Slides                    {slide_count}")
    print("")
    print("[DONE] PPTX generation completed")

def main() -> None:

    """ Run PPTX Generation Workflow """

    parsed_arguments = parse_command_line_arguments()

    # Resolve Generation Paths
    input_markdown_path = resolve_input_markdown_path(parsed_arguments.input_markdown_path)
    output_pptx_path = resolve_output_pptx_path(input_markdown_path, parsed_arguments.output_pptx_path)
    template_pptx_path = resolve_template_pptx_path(parsed_arguments.template_pptx_path)

    # Parse Markdown Slide Deck
    markdown_text = load_markdown_text(input_markdown_path)
    slide_content_list = parse_slide_content_list(markdown_text)

    # Render Presentation
    presentation = create_presentation_from_template(template_pptx_path)
    clear_template_slides(presentation)
    render_slide_content_list(presentation, slide_content_list)
    presentation.save(output_pptx_path)

    # Print Summary
    print_generation_summary(input_markdown_path, output_pptx_path, template_pptx_path, len(slide_content_list))

if __name__ == "__main__":

    main()
